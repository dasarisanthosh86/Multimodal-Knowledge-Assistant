import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import markdown
import os

def extract_from_pdf(file_path):
    """Extract text from PDF using PyMuPDF."""
    try:
        text = ""
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
        return text.strip()
    except Exception as e:
        print(f"❌ Error extracting from PDF: {e}")
        return ""

def extract_from_docx(file_path):
    """Extract text from DOCX files."""
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        print(f"❌ Error extracting from DOCX: {e}")
        return ""

def extract_from_pptx(file_path):
    """Extract text from PowerPoint slides."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"❌ Error extracting from PPTX: {e}")
        return ""

def extract_from_txt(file_path):
    """Extract text from TXT files."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        print(f"❌ Error reading TXT: {e}")
        return ""

def extract_from_md(file_path):
    """Extract text from Markdown and convert to HTML."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()
        html_text = markdown.markdown(md_content)
        return html_text.strip()
    except Exception as e:
        print(f"❌ Error extracting from Markdown: {e}")
        return ""

def extract_text(file_path):
    """Main dispatcher for extracting text from various file types."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            return extract_from_pdf(file_path)
        elif ext == ".docx":
            return extract_from_docx(file_path)
        elif ext == ".pptx":
            return extract_from_pptx(file_path)
        elif ext == ".txt":
            return extract_from_txt(file_path)
        elif ext == ".md":
            return extract_from_md(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        print(f"❌ Extraction failed: {e}")
        return ""
