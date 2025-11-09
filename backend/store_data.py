import os
import uuid
import fitz  # PyMuPDF for PDF
import docx
from pptx import Presentation
from backend.extract_audio import extract_from_audio, extract_from_video
from backend.db import insert_document
from backend.generate_embeddings import create_embeddings  # ✅ NEW: for Gemini embeddings


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from multiple supported file types:
    txt, pdf, docx, pptx, mp3, wav, mp4, mov, avi, png, jpg, jpeg
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    print(f"📂 Extracting from file: {file_path}")

    try:
        # 📝 Plain Text
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        # 📘 PDF
        elif ext == ".pdf":
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text("text")
            doc.close()

        # 📄 Word Document
        elif ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

        # 📊 PowerPoint
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        # 🎵 Audio
        elif ext in [".mp3", ".wav"]:
            text = extract_from_audio(file_path)

        # 🎥 Video
        elif ext in [".mp4", ".mov", ".avi"]:
            text = extract_from_video(file_path)

        # 🖼️ Image (OCR placeholder)
        elif ext in [".png", ".jpg", ".jpeg"]:
            text = "[Image uploaded — OCR not implemented yet.]"

        else:
            print(f"❌ Unsupported file type: {ext}")
            return None

    except Exception as e:
        print(f"❌ Error during extraction from {file_path}: {e}")
        return None

    text = text.strip()
    if not text:
        print("⚠️ No readable text extracted.")
        return None

    print(f"✅ Extraction successful ({len(text)} chars)")
    return text


def process_and_store(file_path: str):
    """
    Extracts text, stores it in MySQL, and creates embeddings.
    Returns: (doc_id, text)
    """
    text = extract_text_from_file(file_path)
    if not text:
        return None, None

    doc_id = str(uuid.uuid4())

    try:
        # ✅ Insert document
        insert_document(doc_id, text)
        print(f"✅ Document stored successfully (ID: {doc_id})")

        # ✅ Create embeddings for Gemini
        create_embeddings(doc_id, text)

        return doc_id, text

    except Exception as e:
        print(f"❌ Database insert or embedding creation failed: {e}")
        return None, None
